from __future__ import annotations

from io import BytesIO
from typing import Any, cast
from xml.sax.saxutils import escape
from zipfile import BadZipFile, ZipFile

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer


DOCUMENT_MIME_TYPES = {
    "docx": (
        "application/vnd.openxmlformats-officedocument."
        "wordprocessingml.document"
    ),
    "xlsx": (
        "application/vnd.openxmlformats-officedocument."
        "spreadsheetml.sheet"
    ),
    "pdf": "application/pdf",
    "pptx": (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.presentation"
    ),
}


def render_document(
    skill_id: str,
    title: str,
    sections: list[tuple[str, list[str]]],
) -> bytes:
    renderers = {
        "docx": _render_docx,
        "xlsx": _render_xlsx,
        "pdf": _render_pdf,
        "pptx": _render_pptx,
    }
    renderer = renderers.get(skill_id)
    if renderer is None:
        raise ValueError(f"未知文档 Skill：{skill_id}")
    content = renderer(title, sections)
    validate_document(skill_id, content)
    return content


def validate_document(skill_id: str, content: bytes) -> None:
    if skill_id == "pdf":
        if not content.startswith(b"%PDF-"):
            raise ValueError("PDF 产物格式校验失败")
        return
    try:
        with ZipFile(BytesIO(content)) as archive:
            names = set(archive.namelist())
    except BadZipFile as exc:
        raise ValueError(f"{skill_id} 产物不是有效压缩包") from exc
    required_member = {
        "docx": "word/document.xml",
        "xlsx": "xl/workbook.xml",
        "pptx": "ppt/presentation.xml",
    }[skill_id]
    if required_member not in names:
        raise ValueError(f"{skill_id} 产物缺少 {required_member}")
    if skill_id == "xlsx":
        workbook = load_workbook(BytesIO(content), read_only=True)
        workbook.close()


def _render_docx(
    title: str,
    sections: list[tuple[str, list[str]]],
) -> bytes:
    document = Document()
    document.add_heading(title, level=0)
    for heading, paragraphs in sections:
        document.add_heading(heading, level=1)
        for text in paragraphs:
            document.add_paragraph(text)
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _render_xlsx(
    title: str,
    sections: list[tuple[str, list[str]]],
) -> bytes:
    workbook = Workbook()
    worksheet = workbook.active
    if worksheet is None:
        raise RuntimeError("XLSX 工作簿没有活动工作表")
    worksheet.title = "报告"
    worksheet.append(["标题", title])
    worksheet.append([])
    worksheet.append(["章节", "内容"])
    for heading, paragraphs in sections:
        for text in paragraphs:
            worksheet.append([heading, text])
    worksheet.column_dimensions["A"].width = 18
    worksheet.column_dimensions["B"].width = 72
    buffer = BytesIO()
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def _render_pdf(
    title: str,
    sections: list[tuple[str, list[str]]],
) -> bytes:
    font_name = "STSong-Light"
    if font_name not in pdfmetrics.getRegisteredFontNames():
        pdfmetrics.registerFont(UnicodeCIDFont(font_name))
    buffer = BytesIO()
    document = SimpleDocTemplate(buffer, pagesize=A4)
    story = [
        Paragraph(
            f'<font name="{font_name}" size="18">{escape(title)}</font>'
        ),
        Spacer(1, 18),
    ]
    for heading, paragraphs in sections:
        story.append(
            Paragraph(
                f'<font name="{font_name}" size="14">'
                f"{escape(heading)}</font>"
            )
        )
        for text in paragraphs:
            story.append(
                Paragraph(
                    f'<font name="{font_name}" size="10">'
                    f"{escape(text)}</font>"
                )
            )
            story.append(Spacer(1, 8))
    document.build(story)
    return buffer.getvalue()


def _render_pptx(
    title: str,
    sections: list[tuple[str, list[str]]],
) -> bytes:
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(
        presentation.slide_layouts[0]
    )
    title_shape = cast(Any, title_slide.shapes.title)
    subtitle_shape = cast(Any, title_slide.placeholders[1])
    title_shape.text = title
    subtitle_shape.text = "由企业智能助理生成"
    for heading, paragraphs in sections:
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[1]
        )
        slide_title = cast(Any, slide.shapes.title)
        slide_body = cast(Any, slide.placeholders[1])
        slide_title.text = heading
        slide_body.text = "\n".join(paragraphs)
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


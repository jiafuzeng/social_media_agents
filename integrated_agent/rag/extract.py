from __future__ import annotations

import logging
from contextlib import contextmanager
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Callable, Iterator

from integrated_agent.rag.clean import clean_document_text
from integrated_agent.rag.models import ExtractError, ExtractOut


@dataclass(frozen=True)
class _Loaded:
    """Reader 产出：纯文本 + 清洗时是否保住换行结构。"""

    text: str
    preserve_structure: bool


def _decode_utf8(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ExtractError(422, "text file is not valid utf-8") from exc


def _read_plain(data: bytes) -> _Loaded:
    return _Loaded(_decode_utf8(data), False)


def _read_markdown(data: bytes) -> _Loaded:
    return _Loaded(_decode_utf8(data), True)


def _read_html(data: bytes) -> _Loaded:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode_utf8(data), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    root = soup.body if soup.body else soup
    return _Loaded(str(root), True)


@contextmanager
def _quiet_pypdf_cmap() -> Iterator[None]:
    """部分中文 PDF 的 ToUnicode CMap 不规范，pypdf 仍能抽出正文但会刷 WARNING。"""

    logger = logging.getLogger("pypdf._cmap")
    previous = logger.level
    logger.setLevel(logging.ERROR)
    try:
        yield
    finally:
        logger.setLevel(previous)


def _read_pdf(data: bytes) -> _Loaded:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ExtractError(500, "pypdf is not installed") from exc
    try:
        reader = PdfReader(BytesIO(data))
    except Exception as exc:
        raise ExtractError(422, "invalid pdf") from exc
    with _quiet_pypdf_cmap():
        pages = [(page.extract_text() or "") for page in reader.pages]
    return _Loaded("\n\n".join(pages), False)


def _read_docx(data: bytes) -> _Loaded:
    try:
        from docx import Document
    except ImportError as exc:
        raise ExtractError(500, "python-docx is not installed") from exc
    try:
        document = Document(BytesIO(data))
    except Exception as exc:
        raise ExtractError(422, "invalid docx") from exc
    parts: list[str] = []
    for paragraph in document.paragraphs:
        body = paragraph.text.strip()
        if body:
            parts.append(body)
    for table in document.tables:
        rows = [
            " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
            for row in table.rows
        ]
        rows = [row for row in rows if row.replace(" | ", "").strip()]
        if rows:
            parts.append("\n".join(rows))
    return _Loaded("\n\n".join(parts), False)


def _read_pptx(data: bytes) -> _Loaded:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ExtractError(500, "python-pptx is not installed") from exc
    try:
        deck = Presentation(BytesIO(data))
    except Exception as exc:
        raise ExtractError(422, "invalid pptx") from exc
    parts: list[str] = []
    for slide in deck.slides:
        for raw in slide.shapes:
            shape: object = raw
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is None:
                continue
            body = str(getattr(text_frame, "text", "")).strip()
            if body:
                parts.append(body)
    return _Loaded("\n\n".join(parts), False)


MAX_UPLOAD_BYTES = 50 * 1024 * 1024

# 对齐 SimpleDirectoryReader.file_extractor：后缀 → Reader。切块仍由用户选 Parser。
FILE_EXTRACTORS: dict[str, Callable[[bytes], _Loaded]] = {
    ".txt": _read_plain,
    ".md": _read_markdown,
    ".markdown": _read_markdown,
    ".html": _read_html,
    ".htm": _read_html,
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
}

_MIME_TO_SUFFIX = {
    "text/plain": ".txt",
    "text/markdown": ".md",
    "text/html": ".html",
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
}


def extract_upload(
    filename: str | None,
    data: bytes,
    content_type: str | None = None,
) -> ExtractOut:
    """按后缀选 Reader 抽出纯文本并清洗。不绑定 embedding、不落冷文件、不入库。"""
    name = (filename or "upload").strip() or "upload"
    mime = (content_type or "").split(";")[0].strip().lower()
    suffix = Path(name).suffix.lower()
    if suffix not in FILE_EXTRACTORS:
        suffix = _MIME_TO_SUFFIX.get(mime, suffix)
    reader = FILE_EXTRACTORS.get(suffix)
    if reader is None:
        raise ExtractError(422, "unsupported file type")
    if len(data) > MAX_UPLOAD_BYTES:
        raise ExtractError(422, "file too large")

    loaded = reader(data)
    text = clean_document_text(loaded.text, preserve_structure=loaded.preserve_structure)
    if not text:
        raise ExtractError(422, "extracted text is empty")
    return ExtractOut(text=text, filename=name, content_type=content_type)

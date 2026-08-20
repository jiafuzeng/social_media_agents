from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path
from unittest.mock import AsyncMock

import pytest
from fastapi.testclient import TestClient
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

from integrated_agent.bootstrap.matrix_service import build_matrix_service
from integrated_agent.config import PROJECT_ROOT
from integrated_agent.rag.chunking import preview_chunks, split_zh_sentences
from integrated_agent.rag.clean import clean_document_text
from integrated_agent.rag.extract import _quiet_pypdf_cmap, extract_upload
from integrated_agent.rag.models import (
    ChunkPreviewError,
    PreviewChunksIn,
)
from tests.fakes import ScriptedMatrixModel, install_scripted_ask
from integrated_agent.transports.http import create_matrix_api

_MARKDOWN = """# 售后

## 退款

七天无理由退款需提供购买凭证。超过七天仅支持质量问题。

## 发票

电子发票在发货后自动开具。

#### 深层备注

这一节不应再按四级标题切开。
"""

_ELEMENT_MD = """# 手册

退款说明如下。

```text
凭证编号
```

| 项目 | 时限 |
| --- | --- |
| 无理由 | 7天 |
"""


def _client(tmp_path: Path, monkeypatch) -> TestClient:
    install_scripted_ask(monkeypatch, ScriptedMatrixModel())
    service = build_matrix_service(
        PROJECT_ROOT,
        identity_root=tmp_path / "identity",
    )
    return TestClient(create_matrix_api(service))


def _auth(client: TestClient, username: str = "owner") -> dict[str, str]:
    created = client.post(
        "/api/users/register",
        json={"username": username, "password": "secret1"},
    )
    assert created.status_code == 201
    return {"Authorization": f"Bearer {created.json()['token']}"}


def test_split_zh_sentences_keeps_terminators() -> None:
    assert split_zh_sentences("你好。世界！行；尾") == ["你好。", "世界！", "行；", "尾"]


def test_split_zh_sentences_does_not_cut_on_pdf_wrap() -> None:
    text = "七天无理由退款需提供购买凭证和完好包装。\n超过七天仅支持质量问题退换。"
    parts = split_zh_sentences(text)
    assert len(parts) == 2
    wrapped = "\n".join(f"第{index}行文字" for index in range(90)) + "。"
    assert len(split_zh_sentences(wrapped)) == 1


def test_clean_document_text_strips_controls_and_wraps() -> None:
    dirty = (
        "\ufeff七天无理由退款需提供\n购买凭证。\r\n\r\n\r\n"
        "超过七天仅支持质量问题。\x00  \n"
    )
    cleaned = clean_document_text(dirty)
    assert cleaned == "七天无理由退款需提供购买凭证。\n\n超过七天仅支持质量问题。"
    assert clean_document_text(cleaned) == cleaned


def test_clean_document_text_keeps_markdown_structure() -> None:
    dirty = "# 售后\n\n\n\n- 七天\n- 质量\n"
    cleaned = clean_document_text(dirty, preserve_structure=True)
    assert cleaned.startswith("# 售后")
    assert "- 七天\n- 质量" in cleaned


@pytest.mark.asyncio
async def test_preview_cleans_file_text_before_parser() -> None:
    wrapped = "\x00".join(["第一句说明文字"] * 40) + "。\r\n第二句。"
    payload = await preview_chunks(
        PreviewChunksIn(text=wrapped, strategy="sentence_window", window_size=1)
    )
    assert len(payload.chunks) == 2
    assert all("\x00" not in chunk.text for chunk in payload.chunks)


@pytest.mark.asyncio
async def test_sentence_window_keeps_wrapped_lines_under_limit() -> None:
    text = "\n".join(f"第{index}行说明文字" for index in range(90)) + "。"
    payload = await preview_chunks(
        PreviewChunksIn(text=text, strategy="sentence_window", window_size=1)
    )
    assert payload.chunks
    assert len(payload.chunks) == 1


@pytest.mark.asyncio
async def test_six_strategies_return_chunks() -> None:
    text = "第一句。第二句。第三句。第四句。"
    for strategy in (
        "sentence",
        "token",
        "markdown",
        "markdown_element",
        "sentence_window",
    ):
        payload = await preview_chunks(
            PreviewChunksIn(text=text, strategy=strategy, chunk_size=512, chunk_overlap=64)
        )
        assert payload.chunks, strategy
        assert payload.notes == ""


@pytest.mark.asyncio
async def test_markdown_chunks_include_header_path() -> None:
    payload = await preview_chunks(PreviewChunksIn(text=_MARKDOWN, strategy="markdown"))
    paths = [chunk.header_path for chunk in payload.chunks if chunk.header_path]
    assert any(path and "售后" in path for path in paths)


@pytest.mark.asyncio
async def test_user_strategy_wins_over_file_suffix() -> None:
    payload = await preview_chunks(
        PreviewChunksIn(text=_MARKDOWN, strategy="sentence", source_suffix=".md")
    )
    assert payload.strategy == "sentence"
    assert "MarkdownNodeParser" not in payload.notes
    assert payload.chunks


@pytest.mark.asyncio
async def test_markdown_card_still_uses_markdown_parser() -> None:
    payload = await preview_chunks(
        PreviewChunksIn(text=_MARKDOWN, strategy="markdown", source_suffix=".md")
    )
    assert "先走" not in payload.notes
    assert any(chunk.header_path for chunk in payload.chunks)


@pytest.mark.asyncio
async def test_html_suffix_does_not_override_sentence_card() -> None:
    html = "<h1>售后</h1><p>七天无理由。</p><p>发票自动开。</p>"
    payload = await preview_chunks(
        PreviewChunksIn(text=html, strategy="sentence", source_suffix="page.html")
    )
    assert "HTMLNodeParser" not in payload.notes
    assert any("七天无理由" in chunk.text for chunk in payload.chunks)


@pytest.mark.asyncio
async def test_markdown_element_types() -> None:
    payload = await preview_chunks(
        PreviewChunksIn(text=_ELEMENT_MD, strategy="markdown_element")
    )
    types = {chunk.element_type for chunk in payload.chunks}
    assert {"title", "text", "code", "table"} <= types


@pytest.mark.asyncio
async def test_sentence_window_has_sentence_and_window() -> None:
    payload = await preview_chunks(
        PreviewChunksIn(
            text="第一句。第二句。第三句。",
            strategy="sentence_window",
            window_size=1,
        )
    )
    assert len(payload.chunks) >= 2
    for chunk in payload.chunks:
        assert chunk.text
        assert chunk.window
        assert chunk.text in chunk.window


@pytest.mark.asyncio
async def test_semantic_embed_stops_after_first_batch_failure(monkeypatch) -> None:
    calls: list[int] = []

    async def _fail(profile_id: str, texts: list[str]) -> list[list[float]]:
        del profile_id
        calls.append(len(texts))
        raise RuntimeError("insufficient API quota")

    monkeypatch.setattr(
        "integrated_agent.rag.embeddings.embed_profile_texts", _fail
    )
    payload = await preview_chunks(
        PreviewChunksIn(
            text="".join(f"句{index}。" for index in range(15)),
            strategy="semantic",
            embedding_profile_id="bge-m3",
        )
    )
    assert calls == [10]
    assert payload.strategy == "sentence"
    assert payload.notes


@pytest.mark.asyncio
async def test_semantic_degrades_when_embed_fails(monkeypatch) -> None:
    async def _fail(profile_id: str, texts: list[str]) -> list[list[float]]:
        del profile_id, texts
        raise RuntimeError("embed unavailable")

    monkeypatch.setattr(
        "integrated_agent.rag.embeddings.embed_profile_texts", _fail
    )
    payload = await preview_chunks(
        PreviewChunksIn(
            text="第一句。第二句。第三句。",
            strategy="semantic",
            embedding_profile_id="bge-m3",
        )
    )
    assert payload.strategy == "sentence"
    assert payload.notes
    assert payload.chunks


@pytest.mark.asyncio
async def test_semantic_uses_step1_embed(monkeypatch) -> None:
    calls: list[list[str]] = []

    async def _embed(profile_id: str, texts: list[str]) -> list[list[float]]:
        assert profile_id == "bge-m3"
        calls.append(list(texts))
        return [[float(index), 1.0, 0.0] for index, _ in enumerate(texts)]

    monkeypatch.setattr(
        "integrated_agent.rag.embeddings.embed_profile_texts", _embed
    )
    payload = await preview_chunks(
        PreviewChunksIn(
            text="退款要凭证。发票自动开。物流三天到。",
            strategy="semantic",
            embedding_profile_id="bge-m3",
        )
    )
    assert calls
    assert payload.notes == ""
    assert payload.strategy == "semantic"
    assert payload.chunks


@pytest.mark.asyncio
async def test_preview_returns_all_chunks() -> None:
    total = 81
    text = "".join(f"句{index}。" for index in range(total))
    payload = await preview_chunks(
        PreviewChunksIn(text=text, strategy="sentence_window", window_size=1)
    )
    assert len(payload.chunks) == total
    assert "只显示前" not in payload.notes


@pytest.mark.asyncio
async def test_unknown_embedding_profile_is_422() -> None:
    with pytest.raises(ChunkPreviewError) as caught:
        await preview_chunks(
            PreviewChunksIn(
                text="第一句。",
                strategy="semantic",
                embedding_profile_id="missing-profile",
            )
        )
    assert caught.value.status == 422


def test_pdf_extract_swallows_broken_cmap_warnings(caplog: pytest.LogCaptureFixture) -> None:
    caplog.set_level(logging.WARNING)
    with _quiet_pypdf_cmap():
        logging.getLogger("pypdf._cmap").warning(
            "Got invalid hex string: Non-hexadecimal digit found (b'CMapName')"
        )
    assert "CMapName" not in caplog.text


def test_extract_txt_and_pdf() -> None:
    extracted = extract_upload("note.md", "售后政策：七天无理由。".encode("utf-8"))
    assert "七天无理由" in extracted.text
    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=A4)
    pdf.drawString(72, 720, "Knowledge extract hello")
    pdf.save()
    from_pdf = extract_upload("doc.pdf", buffer.getvalue(), "application/pdf")
    assert "Knowledge extract hello" in from_pdf.text


def test_extract_markdown_keeps_list_structure() -> None:
    extracted = extract_upload(
        "note.md", "# 售后\n\n- 七天\n- 质量\n".encode("utf-8")
    )
    assert extracted.text.startswith("# 售后")
    assert "- 七天\n- 质量" in extracted.text


def test_extract_html_docx_pptx() -> None:
    from docx import Document
    from pptx import Presentation

    html = extract_upload("page.html", "<p>手册正文。</p>".encode("utf-8"))
    assert "手册正文。" in html.text

    word = Document()
    word.add_paragraph("售后政策：七天无理由。")
    word_buf = BytesIO()
    word.save(word_buf)
    from_docx = extract_upload("note.docx", word_buf.getvalue())
    assert "七天无理由" in from_docx.text

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[0])
    slide.shapes.title.text = "Knowledge extract hello"
    ppt_buf = BytesIO()
    deck.save(ppt_buf)
    from_pptx = extract_upload("deck.pptx", ppt_buf.getvalue())
    assert "Knowledge extract hello" in from_pptx.text


def test_extract_rejects_unsupported_type() -> None:
    from integrated_agent.rag.models import ExtractError

    with pytest.raises(ExtractError) as caught:
        extract_upload("scan.png", b"not-an-image", "image/png")
    assert caught.value.status == 422


def test_extract_rejects_oversize() -> None:
    from integrated_agent.rag.extract import MAX_UPLOAD_BYTES
    from integrated_agent.rag.models import ExtractError

    with pytest.raises(ExtractError) as caught:
        extract_upload("note.txt", b"x" * (MAX_UPLOAD_BYTES + 1))
    assert caught.value.status == 422
    assert str(caught.value) == "file too large"


def test_http_unauthenticated_is_401(tmp_path: Path, monkeypatch) -> None:
    with _client(tmp_path, monkeypatch) as client:
        listed = client.get("/api/kb/chunk-strategies")
        assert listed.status_code == 401
        profiles = client.get("/api/kb/embedding-profiles")
        assert profiles.status_code == 401
        preview = client.post(
            "/api/kb/preview-chunks",
            json={"text": "第一句。", "strategy": "sentence"},
        )
        assert preview.status_code == 401


def test_http_illegal_strategy_is_422_and_all_chunks_are_returned(
    tmp_path: Path, monkeypatch
) -> None:
    with _client(tmp_path, monkeypatch) as client:
        headers = _auth(client)
        illegal = client.post(
            "/api/kb/preview-chunks",
            headers=headers,
            json={"text": "第一句。", "strategy": "chars"},
        )
        assert illegal.status_code == 422
        total = 81
        too_many = client.post(
            "/api/kb/preview-chunks",
            headers=headers,
            json={
                "text": "".join(f"句{index}。" for index in range(total)),
                "strategy": "sentence_window",
                "window_size": 1,
            },
        )
        assert too_many.status_code == 200
        body = too_many.json()
        assert len(body["chunks"]) == total
        assert "只显示前" not in (body.get("notes") or "")


def test_http_preview_does_not_write_record_store(
    tmp_path: Path, monkeypatch
) -> None:
    from integrated_agent.runtimes.matrix.rag import kb_store as kb_store_mod

    blocked = AsyncMock(side_effect=AssertionError("preview must not put"))
    for store in kb_store_mod.kb_store.values():
        monkeypatch.setattr(store, "put", blocked)

    with _client(tmp_path, monkeypatch) as client:
        headers = _auth(client)
        listed = client.get("/api/kb/chunk-strategies", headers=headers)
        assert listed.status_code == 200
        assert listed.json()["default"] == "sentence"
        assert {item["id"] for item in listed.json()["strategies"]} == {
            "sentence",
            "token",
            "markdown",
            "markdown_element",
            "semantic",
            "sentence_window",
        }
        preview = client.post(
            "/api/kb/preview-chunks",
            headers=headers,
            json={"text": _MARKDOWN, "strategy": "markdown"},
        )
        assert preview.status_code == 200
        body = preview.json()
        assert any(chunk.get("header_path") for chunk in body["chunks"])
        token_header = {"X-User-Token": headers["Authorization"].removeprefix("Bearer ")}
        windowed = client.post(
            "/api/kb/preview-chunks",
            headers=token_header,
            json={
                "text": "第一句。第二句。",
                "strategy": "sentence_window",
            },
        )
        assert windowed.status_code == 200
        first = windowed.json()["chunks"][0]
        assert first["text"]
        assert first["window"]
        extracted = client.post(
            "/api/kb/extract",
            headers=headers,
            files={"file": ("note.txt", "手册正文。".encode("utf-8"), "text/plain")},
            data={"embedding_profile_id": "qwen3"},
        )
        assert extracted.status_code == 200
        body = extracted.json()
        assert body["text"] == "手册正文。"
        assert "embedding_profile_id" not in body
    blocked.assert_not_called()


def test_http_embedding_profiles_are_public_labels(tmp_path: Path, monkeypatch) -> None:
    with _client(tmp_path, monkeypatch) as client:
        assert client.get("/api/kb/embedding-profiles").status_code == 401
        headers = _auth(client)
        payload = client.get("/api/kb/embedding-profiles", headers=headers)
        assert payload.status_code == 200
        body = payload.json()
        assert body["default"] == "bge-m3"
        assert body["profiles"] == ["openai-small", "bge-m3", "qwen3"]
        dumped = str(body)
        assert "api_key" not in dumped
        assert "base_url" not in dumped


def test_matrix_page_includes_kb_workspace(tmp_path: Path, monkeypatch) -> None:
    install_scripted_ask(monkeypatch, ScriptedMatrixModel())
    app = create_matrix_api(
        build_matrix_service(PROJECT_ROOT, identity_root=tmp_path / "identity"),
        static_root=PROJECT_ROOT / "static",
    )
    with TestClient(app) as client:
        page = client.get("/matrix").text
        assert 'id="kbWorkspace"' in page
        assert 'data-workspace="kb"' in page
        assert 'id="kbStrategies"' in page
        assert 'id="kbStrategy"' in page
        assert 'id="kbEmbedding"' in page
        assert page.count('id="kbEmbedding"') == 1
        assert 'id="kbDraftProfile"' in page
        assert "/static/matrix-kb.js" in page
        js = (PROJECT_ROOT / "static" / "matrix-kb.js").read_text(encoding="utf-8")
        assert "resetKbWorkspaceState" in js
        assert "kbStorageKey" in js
        assert ">知识库<" in page or "aria-label=\"知识库\"" in page
        assert 'id="kbStepSource"' in page
        assert 'id="kbDrop"' in page
        assert "选择数据源" in page
        assert "尚未写入知识库" in page
        assert "将按以下配置入库" in page
        assert "召回测试" in page
        assert "召回聊天" in page
        assert 'id="kbViewDocs"' in page
        assert 'id="kbViewChat"' in page
        assert page.count('id="kbViewChat"') == 1
        assert 'id="kbChatThread"' in page
        assert 'id="kbChatPlan"' in page
        assert page.count('id="kbChatPlan"') == 1
        assert "一句一个依据" in page
        assert "[[kb:kN]]" in page
        assert "保存并处理" in page
        assert "删除该文档的旧切片记录和向量" in page
        assert "确认重切" in page
